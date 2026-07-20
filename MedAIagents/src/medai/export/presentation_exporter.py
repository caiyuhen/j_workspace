"""
PowerPoint (.pptx) 文档导出模块 (美化版)
Presentation Exporter Module

支持:
- 科研汇报幻灯片
- 影像征象教学幻灯片
- 生物信息学可视化汇报幻灯片
"""

import os
from typing import Dict, List, Any, Optional
from datetime import datetime

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

from .styles import PptxColors, BRAND_NAME, BRAND_TAGLINE


class BasePresentationExporter:
    """PPT 导出基类 - 美化版"""

    # 幻灯片尺寸 16:9
    SLIDE_WIDTH = Inches(13.333)
    SLIDE_HEIGHT = Inches(7.5)

    def __init__(self):
        if not HAS_PPTX:
            raise ImportError("python-pptx is required. Install: pip install python-pptx")
        self.prs = Presentation()
        self.prs.slide_width = self.SLIDE_WIDTH
        self.prs.slide_height = self.SLIDE_HEIGHT

    def _add_decor_bar(self, slide, top=True, height=Inches(0.08)):
        """添加顶部或底部装饰条"""
        if top:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0),
                self.SLIDE_WIDTH, height
            )
        else:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), self.SLIDE_HEIGHT - height,
                self.SLIDE_WIDTH, height
            )
        shape.fill.solid()
        shape.fill.fore_color.rgb = PptxColors.PRIMARY
        shape.line.fill.background()
        # 将装饰条置于底层
        spTree = slide.shapes._spTree
        sp = shape._element
        spTree.remove(sp)
        spTree.insert(2, sp)

    def _add_brand_footer(self, slide):
        """添加品牌页脚"""
        footer = slide.shapes.add_textbox(
            Inches(0.4), self.SLIDE_HEIGHT - Inches(0.35),
            Inches(4), Inches(0.25)
        )
        tf = footer.text_frame
        p = tf.paragraphs[0]
        p.text = BRAND_NAME
        p.font.size = Pt(8)
        p.font.color.rgb = PptxColors.TEXT_LT

    def _add_slide_number(self, slide, number: int):
        """添加页码"""
        num_box = slide.shapes.add_textbox(
            self.SLIDE_WIDTH - Inches(1.0), self.SLIDE_HEIGHT - Inches(0.35),
            Inches(0.8), Inches(0.25)
        )
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(number)
        p.font.size = Pt(10)
        p.font.color.rgb = PptxColors.TEXT_LT
        p.alignment = PP_ALIGN.RIGHT

    def add_title_slide(self, title: str, subtitle: str = "", 
                        author: str = "", doc_type: str = ""):
        """添加专业标题幻灯片（美化版）"""
        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)

        # 背景渐变效果（用纯色矩形模拟）
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.SLIDE_WIDTH, self.SLIDE_HEIGHT
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = PptxColors.PRIMARY_LT
        bg.line.fill.background()
        # 置于底层
        spTree = slide.shapes._spTree
        sp = bg._element
        spTree.remove(sp)
        spTree.insert(2, sp)

        # 左侧色块装饰
        left_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            Inches(0.15), self.SLIDE_HEIGHT
        )
        left_bar.fill.solid()
        left_bar.fill.fore_color.rgb = PptxColors.PRIMARY
        left_bar.line.fill.background()

        # 品牌标识
        brand = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4), Inches(6), Inches(0.4)
        )
        tf = brand.text_frame
        p = tf.paragraphs[0]
        p.text = BRAND_NAME.upper()
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = PptxColors.PRIMARY

        # 文档类型标签
        if doc_type:
            type_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.0), Inches(4), Inches(0.35)
            )
            tf = type_box.text_frame
            p = tf.paragraphs[0]
            p.text = doc_type
            p.font.size = Pt(11)
            p.font.color.rgb = PptxColors.PRIMARY

        # 主标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.0), Inches(12), Inches(1.2)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = PptxColors.PRIMARY_DK

        # 副标题
        if subtitle:
            sub_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(3.3), Inches(12), Inches(0.8)
            )
            tf = sub_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Pt(20)
            p.font.color.rgb = PptxColors.TEXT_LT

        # 作者
        if author:
            auth_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(4.2), Inches(6), Inches(0.4)
            )
            tf = auth_box.text_frame
            p = tf.paragraphs[0]
            p.text = author
            p.font.size = Pt(14)
            p.font.color.rgb = PptxColors.TEXT

        # 日期
        date_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.7), Inches(4), Inches(0.3)
        )
        tf = date_box.text_frame
        p = tf.paragraphs[0]
        p.text = datetime.now().strftime("%Y-%m-%d")
        p.font.size = Pt(11)
        p.font.color.rgb = PptxColors.TEXT_LT

        # 底部装饰条
        self._add_decor_bar(slide, top=False)
        return slide

    def add_content_slide(self, title: str, content_items: List[str],
                          slide_number: int = 0):
        """添加内容幻灯片（美化版）"""
        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)

        # 白色背景
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.SLIDE_WIDTH, self.SLIDE_HEIGHT
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        bg.line.fill.background()
        spTree = slide.shapes._spTree
        sp = bg._element
        spTree.remove(sp)
        spTree.insert(2, sp)

        # 顶部装饰条
        self._add_decor_bar(slide, top=True)

        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.25), Inches(12), Inches(0.7)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = PptxColors.PRIMARY_DK

        # 标题下划线装饰
        underline = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(0.95),
            Inches(1.5), Inches(0.03)
        )
        underline.fill.solid()
        underline.fill.fore_color.rgb = PptxColors.PRIMARY
        underline.line.fill.background()

        # 内容区域
        content_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.8)
        )
        tf = content_box.text_frame
        tf.word_wrap = True

        for i, item in enumerate(content_items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(16)
            p.font.color.rgb = PptxColors.TEXT
            p.space_after = Pt(10)
            p.level = 0

        # 页脚
        self._add_brand_footer(slide)
        if slide_number > 0:
            self._add_slide_number(slide, slide_number)

        return slide

    def add_two_column_slide(self, title: str, left_title: str,
                             left_items: List[str], right_title: str,
                             right_items: List[str], slide_number: int = 0):
        """添加双栏幻灯片（美化版）"""
        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)

        # 白色背景
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.SLIDE_WIDTH, self.SLIDE_HEIGHT
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        bg.line.fill.background()
        spTree = slide.shapes._spTree
        sp = bg._element
        spTree.remove(sp)
        spTree.insert(2, sp)

        # 顶部装饰条
        self._add_decor_bar(slide, top=True)

        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.25), Inches(12), Inches(0.7)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = PptxColors.PRIMARY_DK

        # 标题下划线
        underline = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(0.95),
            Inches(1.5), Inches(0.03)
        )
        underline.fill.solid()
        underline.fill.fore_color.rgb = PptxColors.PRIMARY
        underline.line.fill.background()

        # 左栏标题
        left_title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.2), Inches(5.8), Inches(0.5)
        )
        tf = left_title_box.text_frame
        p = tf.paragraphs[0]
        p.text = left_title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = PptxColors.PRIMARY

        # 左栏内容
        left_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.7), Inches(5.8), Inches(5.0)
        )
        tf = left_box.text_frame
        tf.word_wrap = True
        for item in left_items:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(14)
            p.font.color.rgb = PptxColors.TEXT
            p.space_after = Pt(8)

        # 右栏标题
        right_title_box = slide.shapes.add_textbox(
            Inches(6.8), Inches(1.2), Inches(5.8), Inches(0.5)
        )
        tf = right_title_box.text_frame
        p = tf.paragraphs[0]
        p.text = right_title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = PptxColors.PRIMARY

        # 右栏内容
        right_box = slide.shapes.add_textbox(
            Inches(6.8), Inches(1.7), Inches(5.8), Inches(5.0)
        )
        tf = right_box.text_frame
        tf.word_wrap = True
        for item in right_items:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(14)
            p.font.color.rgb = PptxColors.TEXT
            p.space_after = Pt(8)

        # 页脚
        self._add_brand_footer(slide)
        if slide_number > 0:
            self._add_slide_number(slide, slide_number)

        return slide

    def add_image_slide(self, title: str, image_path: str,
                        caption: str = "", slide_number: int = 0):
        """添加图片幻灯片"""
        slide_layout = self.prs.slide_layouts[6]  # Blank
        slide = self.prs.slides.add_slide(slide_layout)

        # 白色背景
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0),
            self.SLIDE_WIDTH, self.SLIDE_HEIGHT
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        bg.line.fill.background()
        spTree = slide.shapes._spTree
        sp = bg._element
        spTree.remove(sp)
        spTree.insert(2, sp)

        # 顶部装饰条
        self._add_decor_bar(slide, top=True)

        # 标题
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.25), Inches(12), Inches(0.7)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = PptxColors.PRIMARY_DK

        # 插入图片
        if os.path.exists(image_path):
            slide.shapes.add_picture(
                image_path,
                Inches(1.5), Inches(1.3),
                width=Inches(10.3)
            )

        # 图注
        if caption:
            cap_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.5)
            )
            tf = cap_box.text_frame
            p = tf.paragraphs[0]
            p.text = caption
            p.font.size = Pt(12)
            p.font.italic = True
            p.font.color.rgb = PptxColors.TEXT_LT
            p.alignment = PP_ALIGN.CENTER

        # 页脚
        self._add_brand_footer(slide)
        if slide_number > 0:
            self._add_slide_number(slide, slide_number)

        return slide

    def save(self, file_path: str):
        """保存演示文稿"""
        os.makedirs(os.path.dirname(file_path) or '.', exist_ok=True)
        self.prs.save(file_path)
        return file_path


class ResearchPresentationExporter(BasePresentationExporter):
    """科研汇报 PPT 导出器 - 美化版"""

    def export_research_report(self, report_data: Dict[str, Any],
                               file_path: str) -> str:
        """
        将科研报告导出为 PPT

        Args:
            report_data: 报告数据
            file_path: 输出路径
        """
        # 标题页
        self.add_title_slide(
            title=report_data.get("title", "Research Report"),
            subtitle=report_data.get("subtitle", ""),
            author=report_data.get("author", ""),
            doc_type="科研汇报"
        )

        slides = [
            ("Background", report_data.get("background", [])),
            ("Methods", report_data.get("methods", [])),
            ("Results", report_data.get("results", [])),
            ("Discussion", report_data.get("discussion", [])),
            ("Conclusions", report_data.get("conclusions", [])),
            ("Acknowledgments", report_data.get("acknowledgments", [])),
        ]

        slide_num = 2
        for title, items in slides:
            if items:
                self.add_content_slide(title, items, slide_number=slide_num)
                slide_num += 1

        return self.save(file_path)


class ImagingTeachingExporter(BasePresentationExporter):
    """影像征象教学 PPT 导出器 - 美化版"""

    def export_teaching(self, signs: List[Dict[str, Any]],
                        file_path: str) -> str:
        """
        将影像征象教学内容导出为 PPT

        Args:
            signs: 征象列表，每个征象包含 name, description, modalities,
                   anatomy, diseases, severity
            file_path: 输出路径
        """
        self.add_title_slide(
            title="Common Imaging Signs",
            subtitle="Radiology Teaching Series",
            doc_type="影像教学"
        )

        slide_num = 2
        for sign in signs:
            name = sign.get("name", "Unknown Sign")
            description = sign.get("description", "")
            modalities = sign.get("modalities", [])
            anatomy = sign.get("anatomy", [])
            diseases = sign.get("diseases", [])
            severity = sign.get("severity", "")

            left_items = [
                f"Description: {description}",
                f"Modalities: {', '.join(modalities)}",
                f"Anatomy: {', '.join(anatomy)}",
                f"Severity: {severity}",
            ]
            right_items = [f"{i+1}. {d}" for i, d in enumerate(diseases[:5])]

            self.add_two_column_slide(
                name,
                "Basic Information",
                left_items,
                "Related Diseases",
                right_items,
                slide_number=slide_num
            )
            slide_num += 1

        return self.save(file_path)


class BioinformaticsReportExporter(BasePresentationExporter):
    """生物信息学可视化汇报 PPT 导出器 - 美化版"""

    def export_bioinformatics_report(self, report_data: Dict[str, Any],
                                     file_path: str) -> str:
        """
        将生物信息学分析结果导出为 PPT

        Args:
            report_data: 报告数据
            file_path: 输出路径
        """
        self.add_title_slide(
            title=report_data.get("title", "Bioinformatics Analysis Report"),
            subtitle=report_data.get("subtitle", ""),
            author=report_data.get("author", ""),
            doc_type="生信分析"
        )

        slides = [
            ("Sample Information", report_data.get("sample_info", [])),
            ("Mutation Landscape", report_data.get("mutation_summary", [])),
            ("Pathway Enrichment", report_data.get("pathways", [])),
            ("Survival Analysis", report_data.get("survival", [])),
            ("Conclusions", report_data.get("conclusions", [])),
        ]

        slide_num = 2
        for title, items in slides:
            if items:
                self.add_content_slide(title, items, slide_number=slide_num)
                slide_num += 1

        return self.save(file_path)
