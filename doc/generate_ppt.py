from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 定义医疗主题配色
MEDICAL_COLORS = {
    'primary': RGBColor(0, 102, 204),
    'secondary': RGBColor(0, 162, 232),
    'accent': RGBColor(255, 107, 107),
    'text_dark': RGBColor(51, 51, 51),
    'text_light': RGBColor(128, 128, 128),
    'light_bg': RGBColor(245, 250, 255),
    'white': RGBColor(255, 255, 255),
}

# 创建 PPT
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_medicla_background(slide):
    """添加医疗风格的背景元素"""
    # 左侧蓝色条
    left_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.5), prs.slide_height)
    left_shape.fill.solid()
    left_shape.fill.fore_color.rgb = MEDICAL_COLORS['primary']
    left_shape.line.fill.background()
    
    # 装饰圆形
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), Inches(0.5), Inches(2.5), Inches(2.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = MEDICAL_COLORS['light_bg']
    circle.line.fill.background()

def add_title_slide(prs, title, subtitle):
    """添加标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_medicla_background(slide)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(2.5), Inches(12.7), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.name = 'Microsoft YaHei'
    p.font.color.rgb = MEDICAL_COLORS['primary']
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(0.6), Inches(4.2), Inches(12.7), Inches(1.5))
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.name = 'Microsoft YaHei'
    p.font.color.rgb = MEDICAL_COLORS['text_dark']
    p.alignment = PP_ALIGN.CENTER
    
    # 底部信息
    footer = slide.shapes.add_textbox(Inches(0.6), Inches(6.8), Inches(12.7), Inches(0.5))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = "医疗 AI 平台商业计划书 | 2024"
    p.font.size = Pt(12)
    p.font.color.rgb = MEDICAL_COLORS['text_light']
    p.alignment = PP_ALIGN.RIGHT

def add_section_slide(prs, title, number):
    """添加章节页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_medicla_background(slide)
    
    # 数字圆圈
    num_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), Inches(2.2), Inches(1.5), Inches(1.5))
    num_circle.fill.solid()
    num_circle.fill.fore_color.rgb = MEDICAL_COLORS['primary']
    num_circle.line.fill.background()
    
    num_text = num_circle.text_frame
    num_text.text = str(number)
    num_text.paragraphs[0].font.size = Pt(36)
    num_text.paragraphs[0].font.bold = True
    num_text.paragraphs[0].font.color.rgb = MEDICAL_COLORS['white']
    num_text.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 章节标题
    title_box = slide.shapes.add_textbox(Inches(2.2), Inches(2.5), Inches(11), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.name = 'Microsoft YaHei'
    p.font.color.rgb = MEDICAL_COLORS['primary']
    p.alignment = PP_ALIGN.LEFT

def add_content_slide(prs, title, content_items):
    """添加内容页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_medicla_background(slide)
    
    # 标题栏
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12.7), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.name = 'Microsoft YaHei'
    p.font.color.rgb = MEDICAL_COLORS['primary']
    
    # 分割线
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.2), Inches(12.7), Inches(0.1))
    line.fill.solid()
    line.fill.fore_color.rgb = MEDICAL_COLORS['secondary']
    line.line.fill.background()
    
    # 内容区域 - 卡片式布局
    y_start = Inches(1.5)
    card_width = Inches(6.3)
    card_height = Inches(2.3)
    gap = Inches(0.4)
    
    for i, item in enumerate(content_items):
        col = i % 2
        row = i // 2
        x = Inches(0.6) + col * (card_width + gap)
        y = y_start + row * (card_height + gap)
        
        if isinstance(item, dict):
            # 卡片式布局
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_width, card_height)
            card.fill.solid()
            card.fill.fore_color.rgb = MEDICAL_COLORS['light_bg']
            card.line.color.rgb = MEDICAL_COLORS['primary']
            card.line.width = Pt(2)
            
            tf = card.text_frame
            tf.word_wrap = True
            
            # 标题
            p = tf.paragraphs[0]
            p.text = item.get('title', '')
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.name = 'Microsoft YaHei'
            p.font.color.rgb = MEDICAL_COLORS['primary']
            
            # 内容
            if item.get('content'):
                p2 = tf.add_paragraph()
                p2.text = item['content']
                p2.font.size = Pt(12)
                p2.font.name = 'Microsoft YaHei'
                p2.font.color.rgb = MEDICAL_COLORS['text_dark']
                p2.space_after = Pt(0)
        else:
            # 列表项
            box = slide.shapes.add_textbox(Inches(0.6), y, Inches(12.1), Inches(1))
            tf = box.text_frame
            p = tf.paragraphs[0]
            p.text = "• " + item if not item.startswith("• ") else item
            p.font.size = Pt(14)
            p.font.name = 'Microsoft YaHei'
            p.font.color.rgb = MEDICAL_COLORS['text_dark']
            p.space_after = Pt(8)

# 创建幻灯片
add_title_slide(prs, "医疗 AI 平台商业计划书", "EDC × 数字孪生 × 医疗大模型\n助力医疗行业智能化转型\n\n2024 年")

add_section_slide(prs, "01 公司定位与产品矩阵", 1)
add_content_slide(prs, "公司定位与产品矩阵", [
    {
        'title': '使命愿景',
        'content': '使命：用 AI 技术赋能医疗，提升诊疗效率\n愿景：成为全球领先的医疗 AI 平台'
    },
    {
        'title': '核心价值',
        'content': '✓ 提升临床效率 30%+\n✓ 降低医疗成本 25%+\n✓ 改善患者体验'
    },
    {
        'title': 'EDC 系统',
        'content': '电子数据采集系统\n• 临床试验数字化\n• 数据实时质控\n• 合规性保障'
    },
    {
        'title': '患者健康管理系统',
        'content': '• 全程健康管理\n• 慢病监测预警\n• 医患智能互动'
    }
])

add_section_slide(prs, "02 核心技术优势", 2)
add_content_slide(prs, "核心技术优势", [
    {
        'title': '数字孪生患者',
        'content': '真实患者数据的数字化映射\n• 个性化治疗方案模拟\n• 治疗效果预测'
    },
    {
        'title': '数字孪生医生',
        'content': '医生经验的数字化沉淀\n• 辅助诊断决策\n• 知识传承培训'
    },
    {
        'title': '强化学习医疗大模型',
        'content': '基于 Qwen 基底 + RLHF\n• 医疗专业知识 98% 准确率\n• 持续自我优化能力'
    },
    {
        'title': '技术壁垒',
        'content': '✓ 10 万 + 医疗专家标注数据\n✓ 50+ 医疗场景深度优化\n✓ 国家专利 15 项'
    }
])

add_section_slide(prs, "03 快速变现策略", 3)
add_content_slide(prs, "快速变现策略 (0-18 个月)", [
    "阶段 1 (0-6 个月): EDC 系统 + 患者管理系统\n  • 目标收入：500-1000 万\n  • 聚焦 niche 领域 (罕见病、精准医疗)\n  • 免费试用 + 效果付费降低门槛\n\n阶段 2 (6-12 个月): 数字孪生患者 + 医疗大模型\n  • 目标收入：1500-2500 万\n  • 慢病管理先行\n  • 政府购买服务 + 医院合作\n\n阶段 3 (12-18 个月): 数字孪生医生 + 中国移动合作\n  • 目标收入：1000-2000 万\n  • 药企联合研发分成\n  • 中国移动渠道推广"
])

add_section_slide(prs, "04 中国移动合作方案", 4)
add_content_slide(prs, "中国移动合作模式", [
    {
        'title': '模式 1: 联合解决方案 (60%)',
        'content': '• 中国移动：网络 + 算力 + 渠道\n• 我们：AI 平台 + 应用\n• 收益分成：50:50\n• 优势：快速落地，覆盖 31 省'
    },
    {
        'title': '模式 2: 项目制合作 (30%)',
        'content': '• 按项目收费\n• 中国移动作为总包商\n• 利润率：40-50%\n• 适用：智慧医院、区域医疗'
    },
    {
        'title': '模式 3: 合资公司 (10%)',
        'content': '• 双方共同出资成立合资公司\n• 中国移动：49% 股份\n• 我们：51% 股份 + 技术运营\n• 绑定长期利益'
    },
    {
        'title': '预期收益',
        'content': '首年：1000-2000 万\n第二年：3000-5000 万\n第三年：1-2 亿元'
    }
])

add_section_slide(prs, "05 长期战略规划", 5)
add_content_slide(prs, "长期战略规划 (18-60 个月)", [
    "18-36 个月：生态构建期\n  • 建立开发者平台\n  • 开放 API 接口\n  • 引入第三方应用\n  • A 轮融资 (5000 万 -1 亿)\n\n36-48 个月：平台化期\n  • 构建医疗 AI 生态\n  • 跨区域复制\n  • 数据网络效应\n  • B 轮融资 (2-3 亿)\n\n48-60 个月：全球化期\n  • 拓展东南亚市场\n  • 国际化团队\n  • IPO 或并购退出"
])

add_section_slide(prs, "06 财务预测", 6)
add_content_slide(prs, "财务预测 (2024-2027)", [
    "收入预测 (万元)\n  • 2024 年：2000-3000\n  • 2025 年：8000-12000\n  • 2026 年：25000-35000\n  • 2027 年：50000-70000\n\n中国移动合作收益\n  • 2024 年：500-1000 万\n  • 2025 年：2000-3000 万\n  • 2026 年：5000-8000 万\n  • 2027 年：1-2 亿元\n\n盈亏平衡点\n  • 2025 Q3 实现单月盈亏平衡\n  • 2026 Q1 实现年度盈利"
])

add_section_slide(prs, "07 融资规划", 7)
add_content_slide(prs, "融资规划", [
    {
        'title': '天使轮 (2024 Q4)',
        'content': '融资金额：1000-1500 万\n出让股权：10-15%\n用途：产品完善、团队扩充'
    },
    {
        'title': 'A 轮 (2025 Q3)',
        'content': '融资金额：5000-10000 万\n出让股权：15-20%\n用途：市场拓展、生态建设'
    },
    {
        'title': 'B 轮 (2026 Q4)',
        'content': '融资金额：2-3 亿\n出让股权：10-15%\n用途：全球化、并购整合'
    }
])

add_section_slide(prs, "08 风险与应对", 8)
add_content_slide(prs, "风险与应对策略", [
    "政策风险\n  • 应对：紧跟国家政策支持，积极参与标准制定\n  • 合规性保障团队，确保产品符合法规\n\n技术风险\n  • 应对：持续研发投入，保持技术领先\n  • 建立技术护城河，申请核心专利\n\n市场风险\n  • 应对：多渠道推广，降低单一客户依赖\n  • 建立品牌影响力，提升客户粘性\n\n人才风险\n  • 应对：建立有竞争力的薪酬体系\n  • 股权激励计划，吸引高端人才\n\n资金风险\n  • 应对：多元化融资渠道\n  • 严格控制现金流，预留 18 个月运营资金\n\n竞争风险\n  • 应对：差异化竞争，聚焦细分领域\n  • 快速建立规模优势"
])

add_section_slide(prs, "09 关键里程碑", 9)
add_content_slide(prs, "关键里程碑 (2024-2027)", [
    "2024 Q4\n  • EDC 系统完成 V2.0 发布\n  • 签约 10 家医院试点\n  • 天使轮融资完成\n  • 中国移动合作框架签署\n\n2025 Q2\n  • 数字孪生患者上线\n  • 覆盖 50 家医院\n  • 月营收突破 500 万\n  • A 轮启动\n\n2025 Q4\n  • 医疗大模型 1.0 发布\n  • 中国移动合作落地 10 省\n  • 年度营收突破 1 亿\n\n2026 Q4\n  • 生态平台开放\n  • 覆盖 500+ 医院\n  • 年度营收突破 5 亿\n  • B 轮融资完成\n\n2027 Q4\n  • 拓展东南亚市场\n  • 年度营收突破 10 亿\n  • 准备 IPO 或并购"
])

add_section_slide(prs, "10 愿景与结语", 10)
add_content_slide(prs, "愿景与承诺", [
    "我们的愿景\n  • 成为全球领先的医疗 AI 平台\n  • 让每一个患者享受高质量医疗服务\n  • 让每一位医生拥有 AI 超级助手\n\n我们的承诺\n  • 持续技术创新\n  • 严格数据隐私保护\n  • 深度服务临床需求\n  • 与生态伙伴共赢发展\n\n感谢聆听!\n期待与您携手共创医疗 AI 新时代\n\n联系方式：business@medical-ai.com\n电话：400-XXX-XXXX"
])

# 保存文件
output_path = r"d:\workspace\doc\Medical_AI_Business_Plan_Premium.pptx"
prs.save(output_path)
print(f"PPT 已生成并保存至：{output_path}")
print(f"共 {len(prs.slides)} 页幻灯片")
