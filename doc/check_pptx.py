from pptx import Presentation
from pptx.util import Inches, Pt
import os

path = r'd:\doc\患者数字孪生平台技术架构.pptx'
p = Presentation(path)
print(f"Slides: {len(p.slides)}")
for i, slide in enumerate(p.slides):
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text.strip():
            texts.append(shape.text.strip()[:30])
    print(f"  Slide {i+1}: {len(slide.shapes)} shapes, sample texts: {texts[:3]}")

print("File size:", os.path.getsize(path), "bytes")
print("OK!")
