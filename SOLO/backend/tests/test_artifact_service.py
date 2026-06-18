from pathlib import Path

from app.services.artifact_service import ArtifactService


def test_create_markdown_artifact_writes_file_and_returns_download_url(tmp_path):
    service = ArtifactService(base_dir=tmp_path)

    artifact = service.create_markdown_artifact(
        user_id="user-1",
        conversation_id="conv-1",
        task_id="task-1",
        title="糖尿病研究方案",
        content="# 糖尿病研究方案\n\n正文内容",
    )

    artifact_path = Path(artifact["path"])
    assert artifact_path.exists()
    assert artifact_path.read_text(encoding="utf-8") == "# 糖尿病研究方案\n\n正文内容"
    assert artifact["filename"].endswith(".md")
    assert artifact["download_url"] == f"/api/v1/artifacts/{artifact['artifact_id']}/download"
    assert artifact["artifact_id"] in artifact["download_url"]
    assert service.get_owned_artifact(artifact["artifact_id"], "user-1")["task_id"] == "task-1"


def test_create_office_artifacts_for_docx_xlsx_pptx(tmp_path):
    service = ArtifactService(base_dir=tmp_path)

    for fmt in ["docx", "xlsx", "pptx"]:
        artifact = service.create_artifact(
            user_id="user-1",
            conversation_id="conv-1",
            task_id=f"task-{fmt}",
            title=f"测试{fmt}交付物",
            content="# 标题\n\n## 小节\n\n- 要点一\n- 要点二\n\n| 字段 | 内容 |\n|:---|:---|\n| A | B |",
            artifact_format=fmt,
        )
        artifact_path = Path(artifact["path"])
        assert artifact_path.exists()
        assert artifact_path.suffix == f".{fmt}"
        assert artifact["format"] == fmt
        assert artifact["content_type"]
        assert artifact["download_url"] == f"/api/v1/artifacts/{artifact['artifact_id']}/download"


def test_docx_renders_mermaid_flowchart_as_readable_flow_table(tmp_path):
    from docx import Document

    service = ArtifactService(base_dir=tmp_path)
    content = """## 偏离处理流程

```mermaid
graph TD
A[发现偏离] --> B{偏离评估}
B -->|轻微 | C[记录并培训纠正]
B -->|重大 | D[立即通知医学监查员]
D --> E[风险评估：对安全性/数据影响]
E --> F[通知伦理委员会 (如需)]
F --> G[制定纠正预防措施 (CAPA)]
G --> H[跟踪 CAPA 执行]
```
"""

    artifact = service.create_artifact(
        user_id="user-1",
        conversation_id="conv-1",
        task_id="task-mermaid",
        title="偏离处理流程",
        content=content,
        artifact_format="docx",
    )

    document = Document(artifact["path"])
    all_text = "\n".join(p.text for p in document.paragraphs)
    assert "流程图：偏离处理流程" in all_text
    assert "A[发现偏离] --> B{偏离评估}" not in all_text

    table_text = "\n".join(cell.text for table in document.tables for row in table.rows for cell in row.cells)
    assert "节点" in table_text
    assert "类型" in table_text
    assert "条件" in table_text
    assert "下一步" in table_text
    assert "发现偏离" in table_text
    assert "偏离评估" in table_text
    assert "判断" in table_text
    assert "轻微" in table_text
    assert "记录并培训纠正" in table_text
    assert "重大" in table_text
    assert "立即通知医学监查员" in table_text
    assert "跟踪 CAPA 执行" in table_text


def test_pptx_uses_medical_blue_theme_cover_agenda_sections_and_cards(tmp_path):
    from pptx import Presentation
    from pptx.dml.color import RGBColor

    service = ArtifactService(base_dir=tmp_path)
    content = """# 研究背景

GLP-1 受体激动剂可改善血糖控制。
本方案关注 CGM 与肠道菌群分层。

# 研究设计

随机、开放标签、平行对照。
主要终点为 HbA1c 较基线变化。

# 数据管理计划

建立电子数据采集系统。
设置逻辑核查与稽查追踪。
"""

    artifact = service.create_artifact(
        user_id="user-1",
        conversation_id="conv-1",
        task_id="task-themed-pptx",
        title="GLP-1 RCT 方案",
        content=content,
        artifact_format="pptx",
    )

    prs = Presentation(artifact["path"])
    slide_texts = [
        "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))
        for slide in prs.slides
    ]
    all_text = "\n".join(slide_texts)

    assert "医学研究方案" in slide_texts[0]
    assert any("目录" in text for text in slide_texts)
    assert any("01" in text and "研究背景" in text for text in slide_texts)
    assert any("02" in text and "研究设计" in text for text in slide_texts)
    assert any("GLP-1 受体激动剂可改善血糖控制" in text for text in slide_texts)
    assert len(prs.slides) >= 7

    themed_fills = []
    for slide in prs.slides:
        for shape in slide.shapes:
            fill = getattr(shape, "fill", None)
            if fill:
                try:
                    if fill.fore_color.type:
                        themed_fills.append(fill.fore_color.rgb)
                except (AttributeError, TypeError):
                    pass
    assert RGBColor(31, 78, 121) in themed_fills
    assert RGBColor(242, 247, 252) in themed_fills


def test_pptx_renders_mermaid_flowchart_and_markdown_table(tmp_path):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    service = ArtifactService(base_dir=tmp_path)
    content = """## 偏离处理流程

```mermaid
graph TD
A[发现偏离] --> B{偏离评估}
B -->|轻微 | C[记录并培训纠正]
B -->|重大 | D[立即通知医学监查员]
```

## 访视计划

| 阶段 | 检查 |
| --- | --- |
| 筛选 | HbA1c |
| 随访 | 不良事件 |
"""

    artifact = service.create_artifact(
        user_id="user-1",
        conversation_id="conv-1",
        task_id="task-pptx",
        title="RCT方案",
        content=content,
        artifact_format="pptx",
    )

    prs = Presentation(artifact["path"])
    all_text = "\n".join(
        shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")
    )
    assert "A[发现偏离] --> B{偏离评估}" not in all_text
    assert "发现偏离" in all_text
    assert "偏离评估" in all_text
    assert "轻微" in all_text
    assert "记录并培训纠正" in all_text
    assert "重大" in all_text
    assert "立即通知医学监查员" in all_text

    tables = [shape.table for slide in prs.slides for shape in slide.shapes if shape.has_table]
    assert tables, "PPTX 应将 Markdown 表格渲染为真实 PPT 表格"
    table_text = "\n".join(cell.text for table in tables for row in table.rows for cell in row.cells)
    assert "阶段" in table_text
    assert "检查" in table_text
    assert "筛选" in table_text
    assert "HbA1c" in table_text

    connector_count = sum(
        1 for slide in prs.slides for shape in slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.LINE
    )
    assert connector_count >= 2


def test_create_artifact_rejects_unsupported_format(tmp_path):
    service = ArtifactService(base_dir=tmp_path)

    try:
        service.create_artifact(
            user_id="user-1",
            conversation_id="conv-1",
            task_id="task-1",
            title="报告",
            content="内容",
            artifact_format="pdf",
        )
    except ValueError as exc:
        assert "不支持的交付物格式" in str(exc)
    else:
        raise AssertionError("不支持的格式应该被拒绝")


def test_get_owned_artifact_rejects_other_user(tmp_path):
    service = ArtifactService(base_dir=tmp_path)
    artifact = service.create_markdown_artifact(
        user_id="user-1",
        conversation_id="conv-1",
        task_id="task-1",
        title="报告",
        content="内容",
    )

    try:
        service.get_owned_artifact(artifact["artifact_id"], "user-2")
    except PermissionError as exc:
        assert "无权访问" in str(exc)
    else:
        raise AssertionError("其他用户不应访问交付物")
