"""交付物生成与下载记录服务。"""
from __future__ import annotations

import json
import re
import uuid
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Optional

from app.services.artifact_content_cleaner import clean_artifact_content


CONTENT_TYPES = {
    "md": "text/markdown; charset=utf-8",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}


class ArtifactService:
    """管理任务交付物文件。"""

    def __init__(self, base_dir: Optional[Path] = None):
        self.base_dir = Path(base_dir or "deliverables").resolve()
        self.base_dir.mkdir(parents=True, exist_ok=True)
        self.index_path = self.base_dir / "artifacts_index.json"
        self._artifacts: Dict[str, Dict] = self._load_index()

    def _load_index(self) -> Dict[str, Dict]:
        """从磁盘加载交付物索引，避免服务重启后下载链接失效。"""
        if not self.index_path.exists():
            return {}
        try:
            data = json.loads(self.index_path.read_text(encoding="utf-8"))
            if not isinstance(data, dict):
                return {}
            return data
        except Exception:
            return {}

    def _save_index(self) -> None:
        """保存交付物索引到磁盘。"""
        serializable = {}
        for artifact_id, artifact in self._artifacts.items():
            item = dict(artifact)
            created_at = item.get("created_at")
            if hasattr(created_at, "isoformat"):
                item["created_at"] = created_at.isoformat()
            serializable[artifact_id] = item
        self.index_path.write_text(json.dumps(serializable, ensure_ascii=False, indent=2), encoding="utf-8")

    @staticmethod
    def _safe_title(title: str) -> str:
        title = (title or "交付物").strip()[:60]
        title = re.sub(r"[\\/:*?\"<>|\s]+", "_", title)
        return title.strip("_") or "交付物"

    @staticmethod
    def _plain_text(line: str) -> str:
        line = re.sub(r"^#{1,6}\s*", "", line.strip())
        line = re.sub(r"^[-*+]\s+", "", line)
        line = re.sub(r"^\d+[.)]\s+", "", line)
        line = line.replace("**", "").replace("__", "").replace("`", "")
        return line.strip()

    @staticmethod
    def _is_markdown_table_line(line: str) -> bool:
        raw = (line or "").strip()
        return raw.startswith("|") and raw.endswith("|") and raw.count("|") >= 2

    @staticmethod
    def _is_markdown_table_separator(cells: List[str]) -> bool:
        return bool(cells) and all(re.fullmatch(r":?-{3,}:?", (cell or "").strip()) for cell in cells)

    @staticmethod
    def _parse_markdown_table_row(line: str) -> List[str]:
        return [cell.strip() for cell in (line or "").strip().strip("|").split("|")]

    @classmethod
    def _parse_markdown_table(cls, table_lines: List[str]) -> List[List[str]]:
        rows: List[List[str]] = []
        for line in table_lines:
            cells = cls._parse_markdown_table_row(line)
            if cells and not cls._is_markdown_table_separator(cells):
                rows.append(cells)
        return rows

    @staticmethod
    def _parse_mermaid_node(token: str) -> Dict[str, str]:
        """解析 Mermaid 节点 token，例如 A[发现偏离] / B{偏离评估} / D。"""
        raw = (token or "").strip().rstrip(";")
        match = re.match(r"^([A-Za-z0-9_]+)\s*(?:\[(.*?)\]|\{(.*?)\}|\((.*?)\))?$", raw)
        if not match:
            return {"id": raw, "label": raw, "type": "步骤"}
        node_id = match.group(1)
        square_label = match.group(2)
        decision_label = match.group(3)
        round_label = match.group(4)
        label = square_label or decision_label or round_label or node_id
        node_type = "判断" if decision_label is not None else "步骤"
        return {"id": node_id, "label": label.strip(), "type": node_type}

    @classmethod
    def _parse_mermaid_flowchart(cls, mermaid_lines: List[str]) -> List[List[str]]:
        """把简单 Mermaid flowchart/graph TD 转成 Word 可读流程表。"""
        node_meta: Dict[str, Dict[str, str]] = {}
        edges: List[tuple[str, str, str]] = []
        for raw_line in mermaid_lines:
            line = (raw_line or "").strip()
            if not line or line.startswith(("graph ", "flowchart ", "%%")):
                continue
            edge_match = re.match(r"^(?P<src>.+?)\s*-->\|(?P<cond>.*?)\|\s*(?P<dst>.+?)\s*;?$", line)
            if not edge_match:
                edge_match = re.match(r"^(?P<src>.+?)\s*--\|(?P<cond>.*?)\|\s*-->\s*(?P<dst>.+?)\s*;?$", line)
            if not edge_match:
                edge_match = re.match(r"^(?P<src>.+?)\s*-->\s*(?P<dst>.+?)\s*;?$", line)
            if not edge_match:
                continue
            src = cls._parse_mermaid_node(edge_match.group("src"))
            dst = cls._parse_mermaid_node(edge_match.group("dst"))
            for node in (src, dst):
                existing = node_meta.get(node["id"])
                if existing is None or node["label"] != node["id"]:
                    node_meta[node["id"]] = {**(existing or {}), **node}
            condition = (edge_match.groupdict().get("cond") or "").strip()
            edges.append((src["id"], condition, dst["id"]))

        rows = [["节点", "类型", "条件", "下一步"]]
        for src_id, condition, dst_id in edges:
            src = node_meta.get(src_id, {"label": src_id, "type": "步骤"})
            dst = node_meta.get(dst_id, {"label": dst_id, "type": "步骤"})
            rows.append([
                src.get("label", src_id),
                src.get("type", "步骤"),
                condition,
                dst.get("label", dst_id),
            ])
        return rows if len(rows) > 1 else []

    def _artifact_path(self, user_id: str, title: str, artifact_format: str) -> tuple[str, Path]:
        safe_title = self._safe_title(title)
        filename = f"{datetime.now().strftime('%Y%m%d_%H%M%S')}_{safe_title}.{artifact_format}"
        user_dir = self.base_dir / user_id
        user_dir.mkdir(parents=True, exist_ok=True)
        return filename, user_dir / filename

    def _build_metadata(
        self,
        artifact_id: str,
        user_id: str,
        conversation_id: str,
        task_id: str,
        filename: str,
        path: Path,
        artifact_format: str,
    ) -> Dict:
        artifact = {
            "artifact_id": artifact_id,
            "task_id": task_id,
            "conversation_id": conversation_id,
            "user_id": user_id,
            "filename": filename,
            "content_type": CONTENT_TYPES[artifact_format],
            "format": artifact_format,
            "path": str(path),
            "size": path.stat().st_size,
            "created_at": datetime.now(),
            "download_url": f"/api/v1/artifacts/{artifact_id}/download",
        }
        self._artifacts[artifact_id] = artifact
        self._save_index()
        return artifact

    def create_artifact(
        self,
        user_id: str,
        conversation_id: str,
        task_id: str,
        title: str,
        content: str,
        artifact_format: str = "md",
    ) -> Dict:
        """按指定格式创建交付物。"""
        artifact_format = (artifact_format or "md").lower()
        content = clean_artifact_content(content)
        if artifact_format not in CONTENT_TYPES:
            raise ValueError(f"不支持的交付物格式: {artifact_format}")
        if artifact_format == "md":
            return self.create_markdown_artifact(user_id, conversation_id, task_id, title, content)

        artifact_id = str(uuid.uuid4())
        filename, path = self._artifact_path(user_id, title, artifact_format)
        if artifact_format == "docx":
            self._write_docx(path, title, content)
        elif artifact_format == "xlsx":
            self._write_xlsx(path, title, content)
        elif artifact_format == "pptx":
            self._write_pptx(path, title, content)
        return self._build_metadata(artifact_id, user_id, conversation_id, task_id, filename, path, artifact_format)

    def create_markdown_artifact(
        self,
        user_id: str,
        conversation_id: str,
        task_id: str,
        title: str,
        content: str,
    ) -> Dict:
        """保存 Markdown 交付物并返回可下载元数据。"""
        artifact_id = str(uuid.uuid4())
        filename, path = self._artifact_path(user_id, title, "md")
        path.write_text(content or "", encoding="utf-8")
        return self._build_metadata(artifact_id, user_id, conversation_id, task_id, filename, path, "md")

    def _write_docx(self, path: Path, title: str, content: str) -> None:
        from docx import Document
        from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT
        from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
        from docx.oxml import OxmlElement
        from docx.oxml.ns import qn
        from docx.shared import Inches, Pt

        def set_cell_shading(cell, fill: str) -> None:
            tc_pr = cell._tc.get_or_add_tcPr()
            shd = OxmlElement("w:shd")
            shd.set(qn("w:fill"), fill)
            tc_pr.append(shd)

        def set_cell_width(cell, width_inches: float) -> None:
            cell.width = Inches(width_inches)
            tc_pr = cell._tc.get_or_add_tcPr()
            tc_w = tc_pr.first_child_found_in("w:tcW")
            if tc_w is None:
                tc_w = OxmlElement("w:tcW")
                tc_pr.append(tc_w)
            tc_w.set(qn("w:w"), str(int(width_inches * 1440)))
            tc_w.set(qn("w:type"), "dxa")

        def add_markdown_table(rows: List[List[str]]) -> None:
            if not rows:
                return
            col_count = max(len(row) for row in rows)
            normalized_rows = [row + [""] * (col_count - len(row)) for row in rows]
            table = doc.add_table(rows=len(normalized_rows), cols=col_count)
            table.style = "Table Grid"
            table.autofit = True
            widths = self._docx_table_widths(col_count)
            for row_idx, row in enumerate(normalized_rows):
                tr = table.rows[row_idx]
                tr._tr.get_or_add_trPr().append(OxmlElement("w:cantSplit"))
                for col_idx, value in enumerate(row):
                    cell = tr.cells[col_idx]
                    cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.TOP
                    set_cell_width(cell, widths[col_idx])
                    for para in cell.paragraphs:
                        para.alignment = WD_PARAGRAPH_ALIGNMENT.LEFT
                        para.paragraph_format.space_after = Pt(0)
                    cell.text = self._plain_text(value)
                    for para in cell.paragraphs:
                        for run in para.runs:
                            run.font.name = "Arial"
                            run._element.rPr.rFonts.set(qn("w:eastAsia"), "Microsoft YaHei")
                            run.font.size = Pt(9)
                            if row_idx == 0:
                                run.bold = True
                    if row_idx == 0:
                        set_cell_shading(cell, "D9EAF7")
            doc.add_paragraph("")

        def add_mermaid_flowchart(rows: List[List[str]], flow_title: str) -> None:
            if not rows:
                return
            heading = doc.add_paragraph()
            heading.style = "Intense Quote"
            heading.add_run(f"流程图：{self._plain_text(flow_title) or '未命名流程'}").bold = True
            add_markdown_table(rows)

        doc = Document()
        styles = doc.styles
        styles["Normal"].font.name = "Arial"
        styles["Normal"]._element.rPr.rFonts.set(qn("w:eastAsia"), "Microsoft YaHei")
        styles["Normal"].font.size = Pt(11)
        doc.add_heading(self._plain_text(title), level=1)


        lines = (content or "").splitlines()
        idx = 0
        last_heading = self._plain_text(title)
        while idx < len(lines):
            raw = lines[idx].strip()
            if not raw:
                idx += 1
                continue
            if raw.lower() in {"```mermaid", "~~~mermaid"}:
                fence = raw[:3]
                idx += 1
                mermaid_lines: List[str] = []
                while idx < len(lines) and not lines[idx].strip().startswith(fence):
                    mermaid_lines.append(lines[idx])
                    idx += 1
                if idx < len(lines):
                    idx += 1
                add_mermaid_flowchart(self._parse_mermaid_flowchart(mermaid_lines), last_heading)
                continue
            if self._is_markdown_table_line(raw):
                table_rows: List[List[str]] = []
                while idx < len(lines) and self._is_markdown_table_line(lines[idx].strip()):
                    cells = self._parse_markdown_table_row(lines[idx].strip())
                    if not self._is_markdown_table_separator(cells):
                        table_rows.append(cells)
                    idx += 1
                add_markdown_table(table_rows)
                continue
            if raw.startswith("# "):
                last_heading = self._plain_text(raw)
                doc.add_heading(last_heading, level=1)
            elif raw.startswith("## "):
                last_heading = self._plain_text(raw)
                doc.add_heading(last_heading, level=2)
            elif raw.startswith("### "):
                last_heading = self._plain_text(raw)
                doc.add_heading(last_heading, level=3)
            elif raw.startswith(("- ", "* ")):
                doc.add_paragraph(self._plain_text(raw), style="List Bullet")
            else:
                doc.add_paragraph(self._plain_text(raw))
            idx += 1
        doc.save(path)

    @staticmethod
    def _docx_table_widths(col_count: int) -> List[float]:
        """按列数给 Word 表格分配更适合中文研究方案的列宽。"""
        presets = {
            1: [6.4],
            2: [1.9, 4.5],
            3: [1.6, 2.3, 2.5],
            4: [1.3, 1.7, 1.7, 1.7],
            5: [1.1, 1.3, 1.3, 1.3, 1.4],
        }
        if col_count in presets:
            return presets[col_count]
        width = 6.4 / max(col_count, 1)
        return [width] * col_count

    def _extract_table_rows(self, content: str) -> List[List[str]]:
        rows: List[List[str]] = []
        for line in (content or "").splitlines():
            raw = line.strip()
            if not (raw.startswith("|") and raw.endswith("|")):
                continue
            cells = [cell.strip() for cell in raw.strip("|").split("|")]
            if cells and not all(re.fullmatch(r":?-{3,}:?", cell or "") for cell in cells):
                rows.append(cells)
        return rows

    def _write_xlsx(self, path: Path, title: str, content: str) -> None:
        from openpyxl import Workbook
        from openpyxl.styles import Alignment, Font, PatternFill
        from openpyxl.utils import get_column_letter

        wb = Workbook()
        ws = wb.active
        ws.title = "交付物"
        ws["A1"] = self._plain_text(title)
        ws["A1"].font = Font(bold=True, size=14)
        ws["A1"].alignment = Alignment(horizontal="center")
        ws.merge_cells("A1:D1")

        table_rows = self._extract_table_rows(content)
        if table_rows:
            start_row = 3
            for row_idx, row in enumerate(table_rows, start_row):
                for col_idx, value in enumerate(row, 1):
                    cell = ws.cell(row=row_idx, column=col_idx, value=self._plain_text(value))
                    cell.alignment = Alignment(wrap_text=True, vertical="top")
                    if row_idx == start_row:
                        cell.font = Font(bold=True)
                        cell.fill = PatternFill("solid", fgColor="D9EAF7")
        else:
            row_idx = 3
            for line in (content or "").splitlines():
                text = self._plain_text(line)
                if text:
                    ws.cell(row=row_idx, column=1, value=text)
                    row_idx += 1

        for col_idx in range(1, ws.max_column + 1):
            letter = get_column_letter(col_idx)
            max_len = max(len(str(ws.cell(row=row_idx, column=col_idx).value or "")) for row_idx in range(1, ws.max_row + 1))
            ws.column_dimensions[letter].width = min(max_len + 4, 50)
        wb.save(path)

    def _write_pptx(self, path: Path, title: str, content: str) -> None:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.util import Inches, Pt

        prs = Presentation()
        theme = {
            "primary": RGBColor(31, 78, 121),
            "secondary": RGBColor(91, 155, 213),
            "accent": RGBColor(237, 125, 49),
            "background": RGBColor(247, 249, 251),
            "card": RGBColor(242, 247, 252),
            "text": RGBColor(31, 41, 55),
            "muted": RGBColor(107, 114, 128),
            "white": RGBColor(255, 255, 255),
        }
        blank_layout = prs.slide_layouts[6]

        def add_background(slide, color: RGBColor) -> None:
            bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
            bg.fill.solid()
            bg.fill.fore_color.rgb = color
            bg.line.fill.background()

        def add_cover_slide() -> None:
            slide = prs.slides.add_slide(blank_layout)
            add_background(slide, theme["primary"])
            accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.55), Inches(10), Inches(0.25))
            accent.fill.solid()
            accent.fill.fore_color.rgb = theme["accent"]
            accent.line.fill.background()
            deco = slide.shapes.add_shape(MSO_SHAPE.ARC, Inches(7.2), Inches(0.55), Inches(2.4), Inches(2.4))
            deco.line.color.rgb = theme["secondary"]
            deco.line.width = Pt(3)
            title_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.55), Inches(8.6), Inches(1.4))
            p = title_box.text_frame.paragraphs[0]
            p.text = self._plain_text(title)
            p.font.size = Pt(34)
            p.font.bold = True
            p.font.color.rgb = theme["white"]
            subtitle = slide.shapes.add_textbox(Inches(0.75), Inches(3.15), Inches(6.6), Inches(0.5))
            p = subtitle.text_frame.paragraphs[0]
            p.text = "医学研究方案"
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(219, 234, 254)
            tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(4.15), Inches(2.25), Inches(0.45))
            tag.fill.solid()
            tag.fill.fore_color.rgb = theme["accent"]
            tag.line.fill.background()
            p = tag.text_frame.paragraphs[0]
            p.text = "Clinical Research Deck"
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = theme["white"]
            p.alignment = PP_ALIGN.CENTER
            tag.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        add_cover_slide()

        def add_title(slide, slide_title: str) -> None:
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.16), Inches(7.5))
            bar.fill.solid()
            bar.fill.fore_color.rgb = theme["primary"]
            bar.line.fill.background()
            box = slide.shapes.add_textbox(Inches(0.45), Inches(0.25), Inches(9.1), Inches(0.45))
            tf = box.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = self._plain_text(slide_title)[:60]
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = RGBColor(31, 78, 121)

        def add_text_slide(slide_title: str, items: List[str]) -> None:
            slide = prs.slides.add_slide(blank_layout)
            add_background(slide, theme["background"])
            add_title(slide, slide_title)
            for idx, item in enumerate(items[:6]):
                row = idx // 2
                col = idx % 2
                x = 0.65 + col * 4.55
                y = 1.05 + row * 1.65
                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(4.1), Inches(1.25))
                card.fill.solid()
                card.fill.fore_color.rgb = theme["card"]
                card.line.color.rgb = RGBColor(214, 226, 238)
                marker = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.18), Inches(y + 0.25), Inches(0.34), Inches(0.34))
                marker.fill.solid()
                marker.fill.fore_color.rgb = theme["secondary"]
                marker.line.fill.background()
                body = slide.shapes.add_textbox(Inches(x + 0.65), Inches(y + 0.2), Inches(3.2), Inches(0.85)).text_frame
                body.word_wrap = True
                p = body.paragraphs[0]
                p.text = self._plain_text(item)[:140]
                p.font.size = Pt(14)
                p.font.color.rgb = theme["text"]

        def add_table_slide(slide_title: str, rows: List[List[str]]) -> None:
            if not rows:
                return
            slide = prs.slides.add_slide(blank_layout)
            add_background(slide, theme["background"])
            add_title(slide, slide_title)
            col_count = max(len(row) for row in rows)
            normalized = [row + [""] * (col_count - len(row)) for row in rows[:8]]
            table_shape = slide.shapes.add_table(
                len(normalized), col_count, Inches(0.55), Inches(1.05), Inches(9.0), Inches(4.9)
            )
            table = table_shape.table
            for row_idx, row in enumerate(normalized):
                for col_idx, value in enumerate(row):
                    cell = table.cell(row_idx, col_idx)
                    cell.text = self._plain_text(value)[:80]
                    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(13 if col_count > 3 else 15)
                        paragraph.alignment = PP_ALIGN.CENTER
                        if row_idx == 0:
                            paragraph.font.bold = True
                            paragraph.font.color.rgb = RGBColor(255, 255, 255)
                    fill = cell.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor(31, 78, 121) if row_idx == 0 else RGBColor(242, 247, 252)

        def add_flow_slide(slide_title: str, rows: List[List[str]]) -> None:
            if len(rows) <= 1:
                return
            slide = prs.slides.add_slide(blank_layout)
            add_background(slide, theme["background"])
            add_title(slide, f"流程图：{slide_title}")
            edges = rows[1:]
            labels: List[str] = []
            types: Dict[str, str] = {}
            for src, node_type, _condition, dst in edges:
                if src not in labels:
                    labels.append(src)
                if dst not in labels:
                    labels.append(dst)
                types[src] = node_type
                types.setdefault(dst, "步骤")

            positions: Dict[str, tuple[float, float]] = {}
            for idx, label in enumerate(labels[:8]):
                x = 0.65 + (idx % 4) * 2.25
                y = 1.25 + (idx // 4) * 2.05
                positions[label] = (x, y)
                shape_type = MSO_SHAPE.DIAMOND if types.get(label) == "判断" else MSO_SHAPE.ROUNDED_RECTANGLE
                node = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(1.85), Inches(0.78))
                node.fill.solid()
                node.fill.fore_color.rgb = RGBColor(232, 241, 250) if shape_type != MSO_SHAPE.DIAMOND else RGBColor(255, 242, 204)
                node.line.color.rgb = RGBColor(31, 78, 121)
                node.text_frame.clear()
                p = node.text_frame.paragraphs[0]
                p.text = self._plain_text(label)[:24]
                p.font.size = Pt(12)
                p.font.bold = True
                p.alignment = PP_ALIGN.CENTER
                node.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

            for src, _node_type, condition, dst in edges:
                if src not in positions or dst not in positions:
                    continue
                sx, sy = positions[src]
                dx, dy = positions[dst]
                line = slide.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT,
                    Inches(sx + 1.85), Inches(sy + 0.39),
                    Inches(dx), Inches(dy + 0.39),
                )
                line.line.color.rgb = RGBColor(91, 91, 91)
                line.line.width = Pt(1.4)
                if condition:
                    label_box = slide.shapes.add_textbox(
                        Inches((sx + dx) / 2 + 0.6), Inches((sy + dy) / 2 + 0.15), Inches(1.0), Inches(0.28)
                    )
                    p = label_box.text_frame.paragraphs[0]
                    p.text = self._plain_text(condition)[:16]
                    p.font.size = Pt(10)
                    p.font.color.rgb = RGBColor(192, 80, 77)
                    p.alignment = PP_ALIGN.CENTER

        def add_agenda_slide(headings: List[str]) -> None:
            if not headings:
                return
            slide = prs.slides.add_slide(blank_layout)
            add_background(slide, theme["background"])
            add_title(slide, "目录")
            for agenda_idx, heading in enumerate(headings[:8], 1):
                y = 1.0 + (agenda_idx - 1) * 0.58
                num = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(y), Inches(0.55), Inches(0.36))
                num.fill.solid()
                num.fill.fore_color.rgb = theme["primary"]
                num.line.fill.background()
                p = num.text_frame.paragraphs[0]
                p.text = f"{agenda_idx:02d}"
                p.font.size = Pt(11)
                p.font.bold = True
                p.font.color.rgb = theme["white"]
                p.alignment = PP_ALIGN.CENTER
                num.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                box = slide.shapes.add_textbox(Inches(1.45), Inches(y - 0.02), Inches(7.4), Inches(0.38))
                p = box.text_frame.paragraphs[0]
                p.text = self._plain_text(heading)[:50]
                p.font.size = Pt(16)
                p.font.color.rgb = theme["text"]

        def add_section_slide(section_no: int, section_title: str) -> None:
            slide = prs.slides.add_slide(blank_layout)
            add_background(slide, theme["primary"])
            num = slide.shapes.add_textbox(Inches(0.75), Inches(1.0), Inches(1.8), Inches(0.9))
            p = num.text_frame.paragraphs[0]
            p.text = f"{section_no:02d}"
            p.font.size = Pt(42)
            p.font.bold = True
            p.font.color.rgb = theme["accent"]
            title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.25), Inches(8.4), Inches(1.0))
            p = title_box.text_frame.paragraphs[0]
            p.text = self._plain_text(section_title)[:60]
            p.font.size = Pt(30)
            p.font.bold = True
            p.font.color.rgb = theme["white"]
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.82), Inches(3.45), Inches(2.6), Inches(0.08))
            line.fill.solid()
            line.fill.fore_color.rgb = theme["accent"]
            line.line.fill.background()

        lines = (content or "").splitlines()
        headings = [self._plain_text(line.strip()) for line in lines if line.strip().startswith("#")]
        add_agenda_slide(headings)
        section_count = 0
        idx = 0
        last_heading = self._plain_text(title)
        pending_text: List[str] = []

        def flush_text() -> None:
            nonlocal pending_text
            if pending_text:
                chunks = [pending_text[i:i + 7] for i in range(0, len(pending_text), 7)]
                for part_idx, chunk in enumerate(chunks[:8], 1):
                    add_text_slide(last_heading if part_idx == 1 else f"{last_heading}（续）", chunk)
                pending_text = []

        while idx < len(lines):
            raw = lines[idx].strip()
            if not raw:
                idx += 1
                continue
            if raw.startswith("#"):
                flush_text()
                last_heading = self._plain_text(raw)
                section_count += 1
                add_section_slide(section_count, last_heading)
                idx += 1
                continue
            if raw.lower() in {"```mermaid", "~~~mermaid"}:
                flush_text()
                fence = raw[:3]
                idx += 1
                mermaid_lines: List[str] = []
                while idx < len(lines) and not lines[idx].strip().startswith(fence):
                    mermaid_lines.append(lines[idx])
                    idx += 1
                if idx < len(lines):
                    idx += 1
                add_flow_slide(last_heading, self._parse_mermaid_flowchart(mermaid_lines))
                continue
            if self._is_markdown_table_line(raw):
                flush_text()
                table_lines = []
                while idx < len(lines) and self._is_markdown_table_line(lines[idx].strip()):
                    table_lines.append(lines[idx].strip())
                    idx += 1
                add_table_slide(last_heading, self._parse_markdown_table(table_lines))
                continue
            pending_text.append(raw)
            idx += 1
        flush_text()
        prs.save(path)

    def get_owned_artifact(self, artifact_id: str, user_id: str) -> Dict:
        """获取当前用户拥有的交付物。"""
        artifact = self._artifacts.get(artifact_id)
        if not artifact:
            raise FileNotFoundError("交付物不存在")
        if artifact.get("user_id") != user_id:
            raise PermissionError("无权访问此交付物")
        path = Path(artifact["path"])
        if not path.exists() or not path.is_file():
            raise FileNotFoundError("交付物文件不存在或已被删除")
        return artifact


artifact_service = ArtifactService()
